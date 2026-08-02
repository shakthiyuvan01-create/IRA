"""
ppt_builder.py — IRA generates a PowerPoint (.pptx) from a topic.

Content is written by Gemini (a general task → Gemini), then rendered into
slides with python-pptx and saved to Yuvan's Desktop.
"""
import json
import re
import sys
from datetime import datetime
from pathlib import Path


def _base_dir() -> Path:
    if getattr(sys, "frozen", False):
        return Path(sys.executable).parent
    return Path(__file__).resolve().parent.parent


API_CONFIG_PATH = _base_dir() / "config" / "api_keys.json"


def _get_api_key() -> str:
    with open(API_CONFIG_PATH, "r", encoding="utf-8") as f:
        return json.load(f)["gemini_api_key"]


def _desktop() -> Path:
    try:
        from core.user_paths import desktop_dir
        return desktop_dir()
    except Exception:
        return Path.home() / "Desktop"


def _strip_fences(text: str) -> str:
    text = text.strip()
    text = re.sub(r"^```[a-zA-Z]*\n?", "", text)
    text = re.sub(r"\n?```$", "", text)
    return text.strip()


def _generate_outline(topic: str, n: int) -> dict:
    """Ask Gemini for a JSON slide outline."""
    from google import genai
    from google.genai import types

    client = genai.Client(api_key=_get_api_key())
    prompt = (
        f"Create a {n}-slide presentation outline about: {topic}.\n"
        "Return ONLY valid JSON, no markdown, in exactly this shape:\n"
        '{"title": "...", "subtitle": "...", '
        '"slides": [{"title": "...", "bullets": ["...", "..."]}]}\n'
        "Give 3-5 short bullet points per slide. Keep bullets concise."
    )
    resp = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=prompt,
        config=types.GenerateContentConfig(
            system_instruction="You are a presentation expert. Output valid JSON only."
        ),
    )
    raw = _strip_fences((getattr(resp, "text", "") or "").strip())
    return json.loads(raw)


def _fallback_outline(topic: str, n: int) -> dict:
    return {
        "title": topic.title(),
        "subtitle": "Prepared by IRA",
        "slides": [
            {"title": f"{topic.title()} — Overview", "bullets": ["Key point 1", "Key point 2", "Key point 3"]}
            for _ in range(max(1, n - 1))
        ],
    }


def create_presentation(parameters=None, response=None, player=None, session_memory=None) -> str:
    params = parameters or {}
    topic  = (params.get("topic") or params.get("text") or params.get("title") or "").strip()
    if not topic:
        return "What should the presentation be about, Yuvan?"
    try:
        n = int(params.get("slides", 6))
    except Exception:
        n = 6
    n = max(3, min(n, 15))

    try:
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
    except Exception:
        return "python-pptx isn't installed, Yuvan. Run: pip install python-pptx"

    # 1) content
    try:
        data = _generate_outline(topic, n)
        if not isinstance(data, dict) or "slides" not in data:
            raise ValueError("bad shape")
    except Exception as e:
        print(f"[PPT] outline generation failed ({e}); using fallback.")
        data = _fallback_outline(topic, n)

    BLUE = RGBColor(0x0E, 0x86, 0xC9)   # IRA blue accent

    # 2) build deck
    try:
        prs = Presentation()

        # Title slide
        s = prs.slides.add_slide(prs.slide_layouts[0])
        s.shapes.title.text = data.get("title", topic.title())
        try:
            s.placeholders[1].text = data.get("subtitle", "Prepared by IRA")
        except Exception:
            pass
        try:
            s.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = BLUE
        except Exception:
            pass

        # Content slides
        for slide in data.get("slides", []):
            layout = prs.slide_layouts[1]
            cs     = prs.slides.add_slide(layout)
            cs.shapes.title.text = str(slide.get("title", ""))
            try:
                cs.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = BLUE
            except Exception:
                pass
            bullets = slide.get("bullets", []) or []
            try:
                tf = cs.placeholders[1].text_frame
                tf.clear()
                for i, b in enumerate(bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = str(b)
                    p.font.size = Pt(18)
            except Exception:
                pass

        # 3) save
        safe  = re.sub(r"[^\w\- ]", "", topic).strip().replace(" ", "_")[:40] or "presentation"
        ts    = datetime.now().strftime("%Y%m%d_%H%M%S")
        dest  = _desktop()
        dest.mkdir(parents=True, exist_ok=True)
        path  = dest / f"IRA_{safe}_{ts}.pptx"
        n_slides = len(prs.slides._sldIdLst)
        prs.save(str(path))
        return (f"Your presentation on '{topic}' is ready, Yuvan — "
                f"{n_slides} slides saved to your Desktop: {path}")
    except Exception as e:
        return f"I couldn't build the presentation, Yuvan: {e}"
